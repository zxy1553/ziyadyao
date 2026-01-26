from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# 打开PPT文件
prs = Presentation(r'e:\GitHub\ziyadyao\Interview\Oral_Defense.pptx')

# 获取第二张幻灯片（索引为1）
slide = prs.slides[1]

# 定义替换内容
replacements = {
    '开发效率痛点': [
        '需花费大量时间寻找正确的类方法及API，分析继承关系、接口定义、',
        'UE反射、Lua是否可调用复写等，大部分时间不在业务开发上'
    ],
    '调试难度痛点': [
        '遇到无法本地复现或偶现Bug时仅可依靠日志，Bug定位困难；缺乏有效的',
        '日志信息提取方法，人工排查调用链、分析具体场景代码逻辑难度高'
    ],
    '重复工作痛点': [
        'Task开发等框架相似但细节不同的任务，需不断重复开发、微调内部逻辑'
    ]
}

# 找到所有文本框并记录
text_shapes = []
for shape in slide.shapes:
    if hasattr(shape, "text_frame"):
        text_shapes.append(shape)
        print(f"Shape: {shape.text[:50] if len(shape.text) > 50 else shape.text}")

# 遍历所有形状
for i, shape in enumerate(text_shapes):
    if not hasattr(shape, "text_frame"):
        continue
    
    text = shape.text.strip()
    
    # 检查是否包含待填充标记
    if '【待填充' in text or '长等效率问题】' in text or '具等问题】' in text or '成本高等问题】' in text:
        # 这是要替换的描述文本框
        print(f"找到待替换文本框 {i}: {text[:30]}...")
        
        # 清空文本框
        text_frame = shape.text_frame
        text_frame.clear()  # 清空所有段落
        
        # 根据位置判断是哪个痛点
        # 开发效率痛点
        if '效率' in text or '周期' in text:
            for line in replacements['开发效率痛点']:
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(10.5)
                p.font.name = "NotoSansSC"
                p.font.color.rgb = RGBColor(107, 114, 128)
        
        # 调试难度痛点
        elif '调试' in text or 'Bug' in text or '工具' in text:
            for line in replacements['调试难度痛点']:
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(10.5)
                p.font.name = "NotoSansSC"
                p.font.color.rgb = RGBColor(107, 114, 128)
        
        # 重复工作痛点
        elif '重复' in text or '样板' in text or '复用' in text or '维护' in text:
            for line in replacements['重复工作痛点']:
                p = text_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(10.5)
                p.font.name = "NotoSansSC"
                p.font.color.rgb = RGBColor(107, 114, 128)

# 保存文件
prs.save(r'e:\GitHub\ziyadyao\Interview\Oral_Defense_Updated.pptx')
print("PPT更新完成！文件已保存到: e:\\GitHub\\ziyadyao\\Interview\\Oral_Defense_Updated.pptx")